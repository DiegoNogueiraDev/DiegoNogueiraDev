from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# --- COLORS ---
PURPLE_DARK = RGBColor(0x4A, 0x00, 0x80)
PURPLE_MID = RGBColor(0x7B, 0x2D, 0x8E)
PURPLE_LIGHT = RGBColor(0x9D, 0x4E, 0xDD)
PINK = RGBColor(0xE0, 0x5A, 0x9E)
PINK_LIGHT = RGBColor(0xF0, 0x8C, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_SOFT = RGBColor(0xF8, 0xF4, 0xFF)
GRAY_LIGHT = RGBColor(0xE8, 0xDE, 0xF0)

W = prs.slide_width
H = prs.slide_height

# --- BACKGROUND: gradient via full-slide rectangle ---
bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg_shape.rotation = 0
bg_fill = bg_shape.fill
bg_fill.gradient()
bg_fill.gradient_stops[0].position = 0.0
bg_fill.gradient_stops[0].color.rgb = RGBColor(0x1A, 0x00, 0x33)
bg_fill.gradient_stops[1].position = 1.0
bg_fill.gradient_stops[1].color.rgb = RGBColor(0x3D, 0x0A, 0x5C)
bg_shape.line.fill.background()

# --- DECORATIVE: top-right pink accent circle ---
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(-0.8), Inches(3.5), Inches(3.5))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(0x3A, 0x10, 0x55)  # blended pink into dark bg
circle1.line.fill.background()

# --- DECORATIVE: bottom-left purple accent circle ---
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.2), Inches(5.0), Inches(4.0), Inches(4.0))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x2E, 0x0E, 0x48)  # blended purple into dark bg
circle2.line.fill.background()

# --- THIN TOP ACCENT LINE (pink gradient bar) ---
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = PINK
top_bar.line.fill.background()

# ==========================================
# TITLE
# ==========================================
title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.9))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT

run1 = p.add_run()
run1.text = "Problem-First Thinking"
run1.font.size = Pt(38)
run1.font.bold = True
run1.font.color.rgb = WHITE

run2 = p.add_run()
run2.text = "  with AI"
run2.font.size = Pt(38)
run2.font.bold = True
run2.font.color.rgb = PINK

# --- SUBTITLE ---
sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(9.0), Inches(0.5))
tf2 = sub_box.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.LEFT
run_sub = p2.add_run()
run_sub.text = "A new approach where the problem drives every decision — not the framework, not the trend."
run_sub.font.size = Pt(16)
run_sub.font.color.rgb = GRAY_LIGHT
run_sub.font.italic = True

# ==========================================
# MAIN CONTENT: 3 columns with flow
# ==========================================

def add_card(slide, left, top, width, height, icon, title, body, accent_color):
    """Add a rounded card with icon, title, and body text."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x2A, 0x0A, 0x40)
    card.line.color.rgb = RGBColor(0x5A, 0x2A, 0x7A)
    card.line.width = Pt(1)

    # Icon circle
    icon_size = Inches(0.65)
    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + Inches(0.3),
        top + Inches(0.3),
        icon_size, icon_size
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = accent_color
    icon_circle.line.fill.background()

    # Icon text (emoji-like number)
    icon_tf = icon_circle.text_frame
    icon_tf.word_wrap = False
    icon_p = icon_tf.paragraphs[0]
    icon_p.alignment = PP_ALIGN.CENTER
    icon_run = icon_p.add_run()
    icon_run.text = icon
    icon_run.font.size = Pt(22)
    icon_run.font.bold = True
    icon_run.font.color.rgb = WHITE

    # Title
    title_box = slide.shapes.add_textbox(
        left + Inches(1.1), top + Inches(0.25),
        width - Inches(1.4), Inches(0.45)
    )
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    tr.text = title
    tr.font.size = Pt(18)
    tr.font.bold = True
    tr.font.color.rgb = WHITE

    # Body
    body_box = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(1.05),
        width - Inches(0.6), height - Inches(1.3)
    )
    btf = body_box.text_frame
    btf.word_wrap = True
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.LEFT
    bp.space_after = Pt(6)
    br = bp.add_run()
    br.text = body
    br.font.size = Pt(13)
    br.font.color.rgb = GRAY_LIGHT


col_w = Inches(3.7)
col_h = Inches(2.6)
top_y = Inches(2.2)
gap = Inches(0.35)
start_x = Inches(0.8)

# Card 1: Understand
add_card(slide,
    start_x, top_y, col_w, col_h,
    "1",
    "Understand the Problem",
    "Start with the real problem, not with code.\nBreak it down into its essential parts.\nA PRD becomes a structured execution\ngraph — the problem is now visible.",
    PURPLE_LIGHT
)

# Arrow 1→2
arrow1_x = start_x + col_w + Inches(0.02)
arrow1_box = slide.shapes.add_textbox(arrow1_x, top_y + Inches(1.0), Inches(0.32), Inches(0.5))
atf1 = arrow1_box.text_frame
ap1 = atf1.paragraphs[0]
ap1.alignment = PP_ALIGN.CENTER
ar1 = ap1.add_run()
ar1.text = "→"
ar1.font.size = Pt(28)
ar1.font.color.rgb = PINK
ar1.font.bold = True

# Card 2: Structure
add_card(slide,
    start_x + col_w + gap, top_y, col_w, col_h,
    "2",
    "Structure the Solution",
    "AI follows the graph, phase by phase.\n9 stages from ANALYZE to LISTENING.\nEvery task is traceable, validated,\nand connected to the original problem.",
    PINK
)

# Arrow 2→3
arrow2_x = start_x + 2 * col_w + gap + Inches(0.02)
arrow2_box = slide.shapes.add_textbox(arrow2_x, top_y + Inches(1.0), Inches(0.32), Inches(0.5))
atf2 = arrow2_box.text_frame
ap2 = atf2.paragraphs[0]
ap2.alignment = PP_ALIGN.CENTER
ar2 = ap2.add_run()
ar2.text = "→"
ar2.font.size = Pt(28)
ar2.font.color.rgb = PINK
ar2.font.bold = True

# Card 3: Deliver
add_card(slide,
    start_x + 2 * (col_w + gap), top_y, col_w, col_h,
    "3",
    "Deliver with Confidence",
    "The solution stays aligned with the\nproblem at every step. No drift, no\nguesswork. Humans and AI agents\nwork together through the graph.",
    PURPLE_MID
)

# ==========================================
# BOTTOM: Key insight quote
# ==========================================
quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.15), Inches(10.3), Inches(0.7))
qtf = quote_box.text_frame
qtf.word_wrap = True
qp = qtf.paragraphs[0]
qp.alignment = PP_ALIGN.CENTER
qr = qp.add_run()
qr.text = '"The real architecture is not in the code — it\'s in how deeply you understood the problem."'
qr.font.size = Pt(16)
qr.font.italic = True
qr.font.color.rgb = PINK_LIGHT

# ==========================================
# FOOTER: mcp-graph-workflow reference
# ==========================================
footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(6.0), Inches(0.35))
ftf = footer_box.text_frame
ftf.word_wrap = True
fp = ftf.paragraphs[0]
fp.alignment = PP_ALIGN.LEFT
fr1 = fp.add_run()
fr1.text = "mcp-graph-workflow"
fr1.font.size = Pt(13)
fr1.font.bold = True
fr1.font.color.rgb = PURPLE_LIGHT
fr2 = fp.add_run()
fr2.text = "  —  Problem-first AI development  ·  Local-first  ·  Open Source  ·  Made in Brazil"
fr2.font.size = Pt(11)
fr2.font.color.rgb = GRAY_LIGHT

# --- Author ---
author_box = slide.shapes.add_textbox(Inches(9.5), Inches(6.1), Inches(3.5), Inches(0.35))
atf = author_box.text_frame
atp = atf.paragraphs[0]
atp.alignment = PP_ALIGN.RIGHT
atr = atp.add_run()
atr.text = "Diego Nogueira  ·  @NogueiraDev_"
atr.font.size = Pt(11)
atr.font.color.rgb = RGBColor(0x99, 0x88, 0xAA)

# --- BOTTOM ACCENT LINE ---
bot_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.06), W, Inches(0.06))
bot_bar.fill.solid()
bot_bar.fill.fore_color.rgb = PINK
bot_bar.line.fill.background()

# ==========================================
# SAVE
# ==========================================
output_path = "/home/user/DiegoNogueiraDev/problem-first-thinking-with-ai.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
